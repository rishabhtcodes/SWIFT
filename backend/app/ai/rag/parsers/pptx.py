from __future__ import annotations
from pathlib import Path
import pptx

async def parse_pptx(source: str | Path) -> str:
    prs = pptx.Presentation(str(source))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
